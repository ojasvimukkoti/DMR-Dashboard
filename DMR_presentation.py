import streamlit as st
import pandas as pd
# from openpyxl import load_workbook
from datetime import datetime
import plotly.express as px
import plotly.graph_objects as go
import numpy as np
import plotly.io as pio
import os
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from getpass import getuser


def create_bar_chart(df, x_col, y_col, title, xaxis_title, yaxis_title, log_scale = True, bar_color='blue'):
        
    bar_chart=px.bar(
                        df,
                        x=df.index,
                        y=y_col,
                        title = title,
                        labels={y_col:yaxis_title},
                        # text=df[y_col]*100,
                        color_discrete_map={y_col:bar_color}
        )
    if log_scale:
            bar_chart.update_layout(
                xaxis=dict(tickmode='array',tickvals=df[x_col].tolist()),
                xaxis_title=xaxis_title,
                yaxis_title = yaxis_title,
                yaxis_type = 'log',
                bargap=0.1
            )
    else:
            bar_chart.update_layout(
                xaxis=dict(tickmode='linear'),
                xaxis_title=xaxis_title,
                yaxis_title=yaxis_title
            )

    bar_chart.update_traces(textposition='outside', texttemplate='%{y}',
                                textfont=dict(color='black'), marker=dict(color=bar_color))
        
    bar_chart.update_layout(
            xaxis_tickvals=df.index.tolist(),
            xaxis_ticktext=df['DMR #'].tolist()
        )

    return bar_chart

    #Function that will check if a date is in a specfied valid format
def valid_date(date_string,date_format):
    try:
        pd.to_datetime(date_string, format=date_format)
        return True
    except ValueError:
        return False

class DMRAnalysis:
    def __init__(self, data_df):
        self.data_df = data_df

    def read_and_reorder_dmr(self):
        df = pd.read_excel(self.data_df)
            
        # df['Date'] = pd.to_numeric(df['Date'], errors='coerce')

        df['Date']=pd.to_datetime(df['Date'], errors='coerce') #, origin = '1899-12-30'
        # df['Date'] = df['Date'].dt.strftime('%m/%d/%Y')

        return df


    def generate_presentation(self):

        filtered_data = self.read_and_reorder_dmr()
        unnamed_columns = [col for col in filtered_data.columns if 'Unnamed' in col]
        filtered_data = filtered_data.drop(columns=unnamed_columns)

        #changing the data types of columns to different types that will work
        filtered_data["Product Category"] = filtered_data["Product Category"].astype('category')
        filtered_data["Source of Discrepancy"] = filtered_data["Source of Discrepancy"].astype('category')
        filtered_data['Category'] = filtered_data['Category'].astype('category')
        filtered_data["Disposition"] = filtered_data["Disposition"].astype("category")
        filtered_data['Status'] = filtered_data["Status"].astype('category')
        filtered_data["Originator"] = filtered_data["Originator"].astype('str')
        filtered_data["Owner"] = filtered_data["Owner"].astype('str')
        filtered_data['Customer'] = filtered_data["Customer"].astype("str")
        filtered_data['CA Required?'] = filtered_data["CA Required?"].astype("str")

        # print("Dates")
        # print(filtered_data['Date'])

        #Defined the ecpected date format
        expected_date_format = '%m/%d/%Y'

        #Creating list that will hold the invalid date information
        invalid_date_info = []

        #For-loop that will iterate through the Date column and see if the dates are in the correct format
        for index, date_value in filtered_data['Date'].items():
            if not valid_date(date_value, expected_date_format):
                #stores the index of invalid date and the value 
                invalid_date_info.append(date_value)
        #filtering out the dates that not in the correct format and converting the remaining dates
        filtered_data['Date'] = filtered_data['Date'].apply(lambda x: x if valid_date(x, expected_date_format) else pd.NaT)
        #dropping any rows in the data that have 'NaT' as its values in the Date column
        filtered_data.dropna(subset=['Date'], inplace=True)


        #calculating the Age(current date - date) of all DMRs
        current_date = datetime.now()
        formatted_current_date = current_date.strftime('%m-%d-%Y')
        #adding a new column in the filtered_data DF that stores the Age in years
        filtered_data['Age of DMRs (Years)'] = (pd.to_datetime(formatted_current_date) - pd.to_datetime(filtered_data['Date'])).dt.days / 365.25

        #adding a new column in the filteredt DF that has the age of the DMRs in DAYS
        filtered_data['Age of DMRs (Days)'] = (pd.to_datetime(formatted_current_date)-pd.to_datetime(filtered_data['Date'])).dt.days

        #doing the necessary value replacement for the Status column
        # replace_status_map = {'Closed':'CLOSED', 'Open': 'OPEN', 'OPEN ': 'OPEN'}
        # filtered_data['Status'] = filtered_data['Status'].replace(replace_status_map).str.upper().fillna('OTHER')

        #getting the needed data for the Histogram
        open_DMRs_no_disposition = filtered_data[(filtered_data['Disposition'].isnull()) & (filtered_data['Status'] == 'OPEN')]
        open_DMRs_w_disposition = filtered_data[~(filtered_data['Disposition'].isnull()) & (filtered_data['Status']=='OPEN')]

        # Calculate the number of bins dynamically (square root of the number of data points)
        num_bins = int(np.sqrt(len(filtered_data)))

        #calculating the mean of age in days
        mean_no_disposition = open_DMRs_no_disposition['Age of DMRs (Days)'].mean()
        mean_w_disposition = open_DMRs_w_disposition['Age of DMRs (Days)'].mean()

        # print("\n\n")
        # print(filtered_data)
        # print("\n\n")
        # print(open_DMRs_no_disposition)
        # print("\n", open_DMRs_w_disposition)

        #creating histogram of the OPEN DMRs 
        histogram_DMRs = px.histogram(
                    open_DMRs_no_disposition, x='Age of DMRs (Days)', color='Status',
                    nbins=num_bins,
                    labels = {'Age of DMRs':'Age of DMRs(Days)', 'Count':'Frequency'},
                    title='Histogram of OPEN DMRs Ages',
                    color_discrete_map={'OPEN': 'blue'},
                    barmode="overlay"
                )
            #adding a 2nd histogram for the DMRs w/ disposition
        histogram_DMRs.add_trace(
                    px.histogram(
                        open_DMRs_w_disposition, x='Age of DMRs (Days)', color = 'Status',
                        nbins=num_bins, color_discrete_map={'OPEN': 'red'},
                        barmode="overlay"
                    ).data[0]
                )
            #adjusting the legend
        histogram_DMRs.update_layout(
                    legend=dict(x=1, y=1)
                )
            #adding legend for the DMRs w/ and w/o Dispositions
        histogram_DMRs.for_each_trace(lambda t: t.update(legendgroup=t.name))

        histogram_DMRs.data[0].update(name='Triage') #DMRs w/ no dispostion
        histogram_DMRs.data[1].update(name='Disposition') #DMRs with disposition
                
            #Adding vertical line for mean of DMRs with no disposition
        histogram_DMRs.add_shape(
                    go.layout.Shape(
                        type='line',
                        x0=mean_no_disposition,
                        x1=mean_no_disposition,
                        y0=0,
                        y1=1,
                        xref='x',
                        yref='paper',
                        line=dict(color='blue', width=2)
                    )
                )
            #Adding vertical line for mean of DMRs with disposition
        histogram_DMRs.add_shape(
                    go.layout.Shape(
                        type='line',
                        x0=mean_w_disposition,
                        x1=mean_w_disposition,
                        y0=0,
                        y1=1,
                        xref='x',
                        yref='paper',
                        line=dict(color='red', width=2)
                    )
                )
            #adding arrows that point to respective mean on figure
        histogram_DMRs.add_annotation(
                    x=mean_no_disposition, y =1.02,
                    text=f'Mean: {round(mean_no_disposition,3)}',
                    # showarrow = True,
                    xref='x',
                    yref='paper',
                    font= dict(color='blue', size=10),
                    textangle=5
                )
        histogram_DMRs.add_annotation(
                    x=mean_w_disposition, y=1.02,
                    text = f'Mean: {round(mean_w_disposition, 3)}',
                    # showarrow=True,
                    xref='x',
                    yref = 'paper',
                    font=dict(color='red', size=10),
                    textangle=5
                    
                )
            #update traces to add white lines between the bars
        histogram_DMRs.update_traces(marker_line_color='white', marker_line_width=0.5)

        replace_map = {'Rework':'REWORK', 'rework': 'REWORK', 'scrap': 'SCRAP',
                        'REWORk':'REWORK', 'SCRAP ':'SCRAP', "":'OTHER'}

        filtered_data['Disposition'] = filtered_data['Disposition'].replace(replace_map).str.upper().fillna('OTHER')

        filtered_data_disposition = filtered_data[~((filtered_data['Disposition']=='CANCELLED') | (filtered_data['Disposition']=='OTHER'))]

        countS_disposition = filtered_data_disposition['Disposition'].value_counts().reset_index()
        countS_disposition.columns = ['Disposition', 'Count']

            # Sort the DataFrame in descending order based on the 'Count' column
        countS_disposition = countS_disposition.sort_values(by='Count', ascending=False)


        disposition_bar_chart = px.bar(
                        countS_disposition,
                        x = 'Disposition',
                        y='Count',
                        labels={'Disposition':'Disposition', 'index':'Count'},
                        title = 'Counts of DMRs',
                        color_discrete_sequence=['dodgerblue']
                        )                                    

        disposition_bar_chart.update_traces(text=countS_disposition['Count'], textposition='outside')

            #ordering the data to descending order
        descending_ordered_data = filtered_data.sort_values(by='Age of DMRs (Years)', ascending=False)

            #getting the data for the oldest DMRs that have Status of "OPEN"
        oldest_DMRs_open = descending_ordered_data[descending_ordered_data['Status']=='OPEN']
        oldest_open_DMR = oldest_DMRs_open.head(10).reset_index()
        oldest_OPENDMRs = oldest_open_DMR[['DMR #', 'Part #', 'Originator', 'Disposition', 'Status', 'Age of DMRs (Years)']]
            #getting the data for the oldest DMRs with their day ages
        oldest_DMRS_DAYS = oldest_open_DMR[['DMR #', 'Part #', 'Originator', 'Disposition', 'Status', 'Age of DMRs (Days)']]


            #creating graph of DMR ages in days
        dmr_age_days_chart = create_bar_chart(
                            oldest_DMRS_DAYS,
                            x_col = 'DMR #',
                            y_col = 'Age of DMRs (Days)',
                            title = 'Top 10 Oldest OPEN DMRs by Day Age',
                            xaxis_title='DMR Number',
                            yaxis_title='Age of DMRs (Days)',
                            log_scale=True,
                            bar_color='darkorange'
            )

            #create bar chart of DMR counts by QE (NO DISPOSITION)
        qe_counts = open_DMRs_no_disposition['Owner'].value_counts().reset_index().head(10)
        qe_counts.columns = ['Quality Engineer', 'DMR Count']
            #Creating bar chart that charts the QE's DMR Count
        qe_count_chart = px.bar(qe_counts, x = 'Quality Engineer', y='DMR Count',
                                    title = 'DMR Count by Quality Engineer',
                                    labels = {'Quality Engineer': 'Quality Engineer', 'DMR Count': 'DMR Count'},
                                    color_discrete_sequence=['darkviolet'])
            #updating the traces for the chart so it displays the counts
        qe_count_chart.update_traces(textposition='outside', texttemplate='%{y}',
                                        textfont=dict(color='black'))
                

        stats_no_disposition = {
                'Mean': open_DMRs_no_disposition['Age of DMRs (Days)'].mean(),
                'Median': open_DMRs_no_disposition['Age of DMRs (Days)'].median(),
                'Std Dev': open_DMRs_no_disposition['Age of DMRs (Days)'].std(),
                'Range': open_DMRs_no_disposition['Age of DMRs (Days)'].max() - open_DMRs_no_disposition['Age of DMRs (Years)'].min()
            }
            # Calculate statistics for DMRs with disposition
        stats_w_disposition = {
                    'Mean': open_DMRs_w_disposition['Age of DMRs (Days)'].mean(),
                    'Median': open_DMRs_w_disposition['Age of DMRs (Days)'].median(),
                    'Std Dev': open_DMRs_w_disposition['Age of DMRs (Days)'].std(),
                    'Range': open_DMRs_w_disposition['Age of DMRs (Days)'].max() - open_DMRs_w_disposition['Age of DMRs (Years)'].min()
                }

        df_stats_no_disposition = pd.DataFrame(stats_no_disposition, index=['DMRs with no Disposition'])
        df_stats_w_disposition = pd.DataFrame(stats_w_disposition, index=['DMRs with Disposition'])

            #concatenating the two descriptive stats DF's
        df_combinded_stats = pd.concat([df_stats_no_disposition, df_stats_w_disposition])

        template_path = "APCD_PowerPoint_16x9_Template.pptx"
        presentation = Presentation(template_path)

        currentDate = datetime.now().strftime('%m/%d/%Y')
        #creating title slide for presentation
        title_slide_layout = presentation.slide_layouts[0]
        title_slide = presentation.slides.add_slide(title_slide_layout)
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        title.text = f'APCD Daily DMR Log - {currentDate}'
        subtitle.text = "Analysis done by Ojasvi Mukkoti"

            #list that will store all the slides and then go through them to remove the textboxes
        slide_list = []

        graphs = [disposition_bar_chart, dmr_age_days_chart, qe_count_chart]
        graph_names = [ 'Counts of DMRs Chart', 'Top 10 OPEN DMRS DAYS Chart', 'Top 10 DMR Counts by Quality Engineers']

        #adding the histogram slide
        slide2 = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide_list.append(slide2)
        #adding the histogram to the slide
        pio.write_image(histogram_DMRs, 'histogram_open_DMR_ages.png')
        chart1 = slide2.shapes.add_picture('histogram_open_DMR_ages.png', Inches(0.2), Inches(1.1), width=Inches(11.8), height=Inches(6.3))

        #adding a title to the slide
        title_shape = slide2.shapes.title
        title_shape.text = 'Histogram of Open DMRs Ages Chart'


            #getting the histogram slide index to put the stats table in
        histogram_slide_index = presentation.slides[1]

            # Define the dimensions and position of the table
        left = Inches(10.2)
        top = Inches(5)
        width = Inches(3)
        height = Inches(0)

            # Transpose the DataFrame
        df_combinded_stats_transposed = df_combinded_stats.transpose()
        df_combinded_stats_transposed.insert(0, 'Stat', df_combinded_stats_transposed.index)

        table =histogram_slide_index.shapes.add_table(
                rows=df_combinded_stats_transposed.shape[0]+1,
                cols=df_combinded_stats_transposed.shape[1],
                left=left,
                top=top,
                width=width,
                height=height
            )

        table_cells = table.table
        for i, column in enumerate(df_combinded_stats_transposed.columns):
                # Fill in the column headers and set font size
                header_cell = table_cells.cell(0, i)
                header_cell.text = column
                header_cell.text_frame.paragraphs[0].font.size = Pt(8.6)
                for j, value in enumerate(df_combinded_stats_transposed[column]):
                    if isinstance(value, (int, float)):
                        value_str = f"{value:.2f}"
                    else:
                        value_str = str(value)
                    cell = table_cells.cell(j+1, i)
                    cell.text = value_str
                    cell.text_frame.paragraphs[0].font.size=Pt(12)

        for graph, graph_name in zip(graphs, graph_names):
                #saving the graphs as an image
                image_path = f'{graph_name}.png'
                pio.write_image(graph, image_path)
                #adding slide for the graph
                slide = presentation.slides.add_slide(presentation.slide_layouts[1])
                slide_list.append(slide)
                #adding graph as a picture to slide
                chart = slide.shapes.add_picture(image_path, left=Inches(1),top=Inches(1.1), width=Inches(10.7), height=Inches(6.4))
                #updating title of the slides
                title_shape = slide.shapes.title
                title_shape.text = graph_name


        qe_list = qe_counts['Quality Engineer'].tolist()

        top10_QEs_DF = open_DMRs_no_disposition[open_DMRs_no_disposition['Owner'].isin(qe_list)]

            #there is not 'Part Quantity' does it mean: 'Rejected Qty' or 'Final Rejected Qty'??
        columns_to_display = ['Date', 'DMR #', 'RMA #', 'Part #', 'Part Description', 
                                        'Owner', 'Date Code', 'Vendor']

            # Iterate over each QE in top10_QEs_DF
        for i, qe in enumerate(qe_list):
                # Filter data for the current QE
                filtered_data = top10_QEs_DF[top10_QEs_DF['Owner'] == qe]
                
                #selecting top 10 entris for each column
                top_10_entries = filtered_data[columns_to_display].head(10)
                # Create a new slide for the current QE
                slide = presentation.slides.add_slide(presentation.slide_layouts[1])
                
                slide_list.append(slide)
                
                # Set title for the slide
                title_text = f"Table for {qe} Quality Engineer"
                title_shape = slide.shapes.title
                title_shape.text = title_text

                slide_width = Inches(10)
                slide_height = Inches(7)

                # Define the dimensions and position of the table
                left = Inches(0.4)
                top = Inches(1.1)
                width = Inches(12.6)
                height = Inches(0)

                if (left+width > slide_width) or (top+height >slide_height):
                    font_size = Pt(8)
                else:
                    font_size = Pt(10)

                # Add the table to the slide
                table = slide.shapes.add_table(
                    rows=top_10_entries.shape[0] + 1,
                    cols=top_10_entries.shape[1],
                    left=left,
                    top=top,
                    width=width,
                    height=height
                )

                # Fill in the table with data
                table_cells = table.table
                for col_index, column in enumerate(top_10_entries.columns):
                    # Fill in the column headers and setting the font size
                    header_cell = table_cells.cell(0, col_index)
                    header_cell.text = column
                    header_cell.text_frame.paragraphs[0].font.size = font_size
                    for row_index, value in enumerate(top_10_entries[column]):
                        if pd.notnull(value):
                            if isinstance(value, (int, float)):
                                value_str = f"{value:.2f}"
                            else:
                                value_str = str(value)
                            cell = table_cells.cell(row_index + 1, col_index)
                            cell.text = value_str
                            # Set font size for the text in the cell
                            cell.text_frame.paragraphs[0].font.size = font_size

                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.size=font_size

            # #iterating over the shapes in each slide and removing the big textboxes
        for slide in slide_list:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        #checking is the shape is not a title textbox
                        if not shape==slide.shapes.title:
                            slide.shapes._spTree.remove(shape._element) 
                            
            #saving the presentation to the APCD DMR Log folder
        output_directory = r'P:\Quality\DMR\APCD Daily DMR Log'
        output_filename = f'DMR_Analysis_{formatted_current_date}.pptx'
        full_output_path = os.path.join(output_directory, output_filename).replace('/', '\\')

        presentation.save(full_output_path)

